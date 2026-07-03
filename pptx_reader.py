"""
Extracts the LNG topics directly from the source .pptx so the
Streamlit app always reflects the deck, not a hand-maintained copy.

The deck follows a consistent per-slide template: a bare topic number, an
ALL-CAPS category label, a subtitle sentence, then a body made of ALL-CAPS
sub-headers, bullet paragraphs, numbered items (number shape + title/body
shape pair), stat callouts (short number shape + caption shape pair), and
occasional tables. Some topics span two slides (footer "(1 of 2)"/"(2 of 2)"),
whose bodies are concatenated in slide order.
"""
import re

from pptx import Presentation

_NUM_RE = re.compile(r"^\d{1,2}$")
_STAT_RE = re.compile(r"^[≈~]?[$€]?[\d.,–—-]+%?$")
_FOOTER_RE = re.compile(r"^Q\d+ of \d+")
_PART_SUFFIX_RE = re.compile(r"\s*\(\d+ of \d+\)\s*$")


def _paragraphs(shape):
    return [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]


def _table_to_markdown(table):
    rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    header, *body = rows
    lines = ["| " + " | ".join(header) + " |", "|" + "|".join(["---"] * len(header)) + "|"]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _slide_items(slide):
    items = []
    for shape in slide.shapes:
        if shape.has_table:
            items.append(("table", shape))
        elif shape.has_text_frame and shape.text_frame.text.strip():
            items.append(("text", shape))
    return items


def _parse_body(items, start):
    lines = []
    i, n = start, len(items)
    while i < n:
        kind, shape = items[i]
        if kind == "table":
            lines.append(_table_to_markdown(shape.table))
            i += 1
            continue

        paras = _paragraphs(shape)
        if not paras:
            i += 1
            continue
        if _FOOTER_RE.match(" ".join(paras)):
            i += 1
            continue

        # Numbered item: bare "N" shape followed by a (title, body) shape.
        if len(paras) == 1 and _NUM_RE.match(paras[0]) and i + 1 < n:
            nxt_paras = _paragraphs(items[i + 1][1]) if items[i + 1][0] == "text" else []
            if nxt_paras:
                title, rest = nxt_paras[0], " ".join(nxt_paras[1:]).strip()
                lines.append(f"{paras[0]}. **{title}**" + (f" — {rest}" if rest else ""))
                i += 2
                continue

        # Stat callout: short numeric/currency shape followed by a caption.
        if (
            len(paras) == 1
            and _STAT_RE.match(paras[0])
            and any(ch.isdigit() for ch in paras[0])
            and i + 1 < n
        ):
            nxt_paras = _paragraphs(items[i + 1][1]) if items[i + 1][0] == "text" else []
            if nxt_paras and not nxt_paras[0].isupper():
                lines.append(f"**{paras[0]}** — {' '.join(nxt_paras)}")
                i += 2
                continue

        # ALL-CAPS sub-header.
        if len(paras) == 1 and paras[0].isupper():
            lines.append(f"**{paras[0]}**")
            i += 1
            continue

        # Regular bullet content: one bullet per paragraph.
        for para in paras:
            lines.append(f"- {para}")
        i += 1

    return lines


def load_topics(pptx_path):
    """Returns (intro_markdown, [(title, body_markdown), ...]) in slide order."""
    slides = list(Presentation(pptx_path).slides)

    intro_paras = [p for _, shape in _slide_items(slides[0]) for p in _paragraphs(shape)]
    intro = f"# {intro_paras[0]}\n\n" + "\n\n".join(f"**{intro_paras[1]}**" if i == 0 else t
                                                       for i, t in enumerate(intro_paras[1:]))

    topics, order = {}, []
    for slide in slides[1:]:
        items = _slide_items(slide)
        if not items or items[0][0] != "text":
            continue
        num_paras = _paragraphs(items[0][1])
        if not (len(num_paras) == 1 and _NUM_RE.match(num_paras[0])):
            continue
        topic_no = int(num_paras[0])
        idx = 1

        category = ""
        if idx < len(items) and items[idx][0] == "text":
            p = _paragraphs(items[idx][1])
            if len(p) == 1 and p[0].isupper():
                category, idx = p[0], idx + 1

        subtitle = ""
        if idx < len(items) and items[idx][0] == "text":
            p = _paragraphs(items[idx][1])
            if p:
                subtitle, idx = p[0], idx + 1
        subtitle = _PART_SUFFIX_RE.sub("", subtitle).strip()

        body_lines = _parse_body(items, idx)

        if topic_no not in topics:
            topics[topic_no] = {"title": subtitle, "category": category, "lines": []}
            order.append(topic_no)
        topics[topic_no]["lines"].extend(body_lines)

    sections = []
    for n in order:
        t = topics[n]
        header = f"## Q{n}. {t['title']}"
        if t["category"]:
            header += f"\n\n*{t['category']}*"
        body = "\n\n".join(t["lines"])
        sections.append((f"Q{n}. {t['title']}", f"{header}\n\n{body}"))

    return intro, sections
